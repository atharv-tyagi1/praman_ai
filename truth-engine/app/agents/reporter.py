"""
Agent 7 — Reporter
Generates professional PPT and PDF reports from fact-verification results.
Uses python-pptx for slides and fpdf2 for PDF documentation.
"""

import os
import logging
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from fpdf import FPDF

# Ensure reports directory exists
REPORTS_DIR = "reports"
os.makedirs(REPORTS_DIR, exist_ok=True)

class Reporter:
    @staticmethod
    async def generate_ppt(report_data: dict, session_id: str) -> str:
        """Generates a professional PPT report."""
        try:
            prs = Presentation()
            
            # 1. Title Slide
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = "Fact Verification Analysis Report"
            overall = report_data.get("overall_assessment", {})
            credibility = round(overall.get("overall_credibility", 0) * 100)
            subtitle.text = f"Session: {session_id}\nCredibility Score: {credibility}%\nGenerated on: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
            
            # 2. Executive Summary
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = "Executive Summary"
            body = slide.placeholders[1]
            body.text = overall.get("summary", "No summary available.")
            
            # 3. Claims & Verdicts
            verdicts = report_data.get("verdicts", [])
            for verdict in verdicts:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = f"Claim Analysis: {verdict.get('verdict', 'Unknown')}"
                
                # Content Box
                tf = slide.placeholders[1].text_frame
                tf.text = f"Claim: {verdict.get('claim_text')}\n"
                tf.add_paragraph().text = f"Verdict: {verdict.get('verdict')}"
                tf.add_paragraph().text = f"Confidence: {round(verdict.get('confidence', 0) * 100)}%"
                tf.add_paragraph().text = f"\nExplanation: {verdict.get('explanation')}"
                
                if verdict.get("key_evidence"):
                    p = tf.add_paragraph()
                    p.text = f"\nKey Evidence: {verdict.get('key_evidence')}"
                    p.font.bold = True
            
            output_path = os.path.join(REPORTS_DIR, f"report_{session_id}.pptx")
            prs.save(output_path)
            logging.info(f"PPT generated: {output_path}")
            return output_path
        except Exception as e:
            logging.error(f"Error generating PPT: {e}")
            return ""

    @staticmethod
    async def generate_pdf(report_data: dict, session_id: str) -> str:
        """Generates a structured PDF report."""
        try:
            pdf = FPDF()
            pdf.add_page()
            
            # Header
            pdf.set_font("Helvetica", "B", 24)
            pdf.cell(0, 20, "Praman-AI Verification Report", ln=True, align="C")
            
            pdf.set_font("Helvetica", "", 12)
            pdf.cell(0, 10, f"Session ID: {session_id}", ln=True, align="C")
            pdf.cell(0, 10, f"Date: {datetime.now().strftime('%Y-%m-%d %H:%M')}", ln=True, align="C")
            pdf.ln(10)
            
            # Score Overview
            overall = report_data.get("overall_assessment", {})
            credibility = round(overall.get("overall_credibility", 0) * 100)
            
            pdf.set_font("Helvetica", "B", 16)
            pdf.set_fill_color(240, 240, 240)
            pdf.cell(0, 15, "Credibility Overview", ln=True, fill=True)
            pdf.set_font("Helvetica", "", 12)
            pdf.cell(0, 10, f"Overall Credibility Score: {credibility}%", ln=True)
            pdf.multi_cell(0, 8, f"Summary: {overall.get('summary', 'N/A')}")
            pdf.ln(10)
            
            # Claims
            pdf.set_font("Helvetica", "B", 16)
            pdf.set_fill_color(240, 240, 240)
            pdf.cell(0, 15, "Claim-by-Claim Breakdown", ln=True, fill=True)
            pdf.ln(5)
            
            verdicts = report_data.get("verdicts", [])
            for i, v in enumerate(verdicts):
                pdf.set_font("Helvetica", "B", 12)
                pdf.multi_cell(0, 8, f"Claim {i+1}: {v.get('claim_text')}")
                
                # Verdict Badge (Simple color coding)
                verdict_str = v.get("verdict", "Unknown")
                pdf.set_font("Helvetica", "I", 11)
                pdf.cell(0, 8, f"Verdict: {verdict_str} | Confidence: {round(v.get('confidence', 0) * 100)}%", ln=True)
                
                pdf.set_font("Helvetica", "", 10)
                pdf.multi_cell(0, 6, f"Explanation: {v.get('explanation')}")
                
                if v.get("key_evidence"):
                    pdf.set_font("Helvetica", "B", 10)
                    pdf.multi_cell(0, 6, f"Key Evidence: {v.get('key_evidence')}")
                
                pdf.ln(5)
                pdf.line(10, pdf.get_y(), 200, pdf.get_y())
                pdf.ln(5)
            
            output_path = os.path.join(REPORTS_DIR, f"report_{session_id}.pdf")
            pdf.output(output_path)
            logging.info(f"PDF generated: {output_path}")
            return output_path
        except Exception as e:
            logging.error(f"Error generating PDF: {e}")
            return ""
